"""
PPT导出功能 - 可编辑文本版本
"""
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import os
from PIL import Image, ImageDraw, ImageFont


class PPTExporter:
    """PPT导出器"""

    # 96 DPI: 1px = 914400 / 96 = 9525 EMU
    PIXELS_TO_EMU = 9525
    MAX_PPT_PIXELS = 5000

    def __init__(self, text_bg_color=None, text_bg_alpha=200, slide_size_px=None, allow_upscale=False):
        """
        初始化PPT

        Args:
            text_bg_color: 文本框背景颜色 (R, G, B) 或 None（无背景）
            text_bg_alpha: 文本框背景不透明度 0-255（仅 text_bg_color 不为 None 时生效）
            slide_size_px: 固定幻灯片尺寸 (width_px, height_px)。混合尺寸导出时应预先传入。
            allow_upscale: 当图片小于固定幻灯片尺寸时，是否放大填充。
        """
        self.prs = Presentation()
        self.text_bg_color = text_bg_color  # 例如: (255, 255, 255) 白色
        try:
            self.text_bg_alpha = int(text_bg_alpha)
        except Exception:
            self.text_bg_alpha = 200
        self.dimensions_set = False  # 标记是否已设置尺寸
        self.slide_size_px = self._normalize_slide_size(slide_size_px)
        self.allow_upscale = bool(allow_upscale)

    @classmethod
    def _scale_to_ppt_limit(cls, img_width, img_height):
        img_width = max(1, int(img_width or 1))
        img_height = max(1, int(img_height or 1))
        max_dimension = max(img_width, img_height)
        if max_dimension > cls.MAX_PPT_PIXELS:
            ppt_scale = cls.MAX_PPT_PIXELS / max_dimension
            scaled_width = max(1, int(round(img_width * ppt_scale)))
            scaled_height = max(1, int(round(img_height * ppt_scale)))
        else:
            ppt_scale = 1.0
            scaled_width = img_width
            scaled_height = img_height
        return scaled_width, scaled_height, float(ppt_scale)

    @classmethod
    def presentation_size_for_images(cls, image_paths):
        max_w = 1
        max_h = 1
        for image_path in image_paths or []:
            try:
                with Image.open(image_path) as img:
                    scaled_w, scaled_h, _ = cls._scale_to_ppt_limit(*img.size)
                max_w = max(max_w, scaled_w)
                max_h = max(max_h, scaled_h)
            except Exception:
                continue
        return max_w, max_h

    def _normalize_slide_size(self, slide_size_px):
        if not (isinstance(slide_size_px, (list, tuple)) and len(slide_size_px) == 2):
            return None
        try:
            w = max(1, int(round(float(slide_size_px[0]))))
            h = max(1, int(round(float(slide_size_px[1]))))
        except Exception:
            return None
        return (w, h)

    def _resolve_slide_layout(self, img_width, img_height):
        scaled_width, scaled_height, ppt_scale = self._scale_to_ppt_limit(img_width, img_height)
        slide_w, slide_h = self.slide_size_px or (scaled_width, scaled_height)

        layout_scale = min(float(slide_w) / float(scaled_width), float(slide_h) / float(scaled_height))
        if not self.allow_upscale:
            layout_scale = min(1.0, layout_scale)
        layout_scale = max(layout_scale, 1e-6)

        placed_w = max(1, int(round(float(scaled_width) * layout_scale)))
        placed_h = max(1, int(round(float(scaled_height) * layout_scale)))
        left = max(0, int(round((float(slide_w) - float(placed_w)) / 2.0)))
        top = max(0, int(round((float(slide_h) - float(placed_h)) / 2.0)))

        return {
            "slide_width": int(slide_w),
            "slide_height": int(slide_h),
            "image_width": placed_w,
            "image_height": placed_h,
            "left": left,
            "top": top,
            "ppt_scale": float(ppt_scale),
            "layout_scale": float(layout_scale),
            "content_scale": float(ppt_scale) * float(layout_scale),
        }

    def _ensure_slide_dimensions(self, slide_w, slide_h):
        pixels_to_emu = self.PIXELS_TO_EMU
        if not self.dimensions_set:
            self.prs.slide_width = int(slide_w * pixels_to_emu)
            self.prs.slide_height = int(slide_h * pixels_to_emu)
            self.dimensions_set = True
            return

        current_w = int(round(self.prs.slide_width / pixels_to_emu))
        current_h = int(round(self.prs.slide_height / pixels_to_emu))
        if current_w != int(slide_w) or current_h != int(slide_h):
            raise RuntimeError("PPT slide size is already fixed. Pass a shared slide_size_px before adding slides.")

    def _rect_to_xywh(self, rect, img_w, img_h):
        """
        将 rect 归一化为 (x, y, w, h)。

        支持两种输入：
        - (x, y, w, h)
        - (x1, y1, x2, y2)
        """
        try:
            x, y, a, b = rect
        except Exception:
            return 0, 0, 1, 1

        # 容错：允许轻微越界（例如四舍五入导致的 1-2px 偏差）
        tol = 2

        # 作为 (x, y, w, h) 的合法性判断
        xywh_ok = (
            a > 0 and b > 0
            and x >= 0 and y >= 0
            and (x + a) <= (img_w + tol)
            and (y + b) <= (img_h + tol)
        )

        # 作为 (x1, y1, x2, y2) 的合法性判断
        w2 = a - x
        h2 = b - y
        xyxy_ok = (
            w2 > 0 and h2 > 0
            and x >= 0 and y >= 0
            and a <= (img_w + tol)
            and b <= (img_h + tol)
        )

        # 若只有一种解释合理，则选择它；都合理时默认按 xywh（本项目约定）
        if xyxy_ok and not xywh_ok:
            return x, y, w2, h2

        if xywh_ok:
            return x, y, a, b

        # 都不合理，按 xywh 回退并做基本修正，避免生成异常/巨大文本框
        x = max(0, int(round(x)))
        y = max(0, int(round(y)))
        w = max(1, int(round(a)))
        h = max(1, int(round(b)))
        return x, y, w, h

    def add_image_with_text_boxes(self, image_path: str, text_boxes: list, title: str = ""):
        """
        添加带可编辑文本框的图片页（按图片原始大小）

        Args:
            image_path: 图片路径
            text_boxes: 文本框列表，每项包含 rect, text, confidence
            title: 页面标题（不使用）
        """
        # 创建空白页
        blank_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(blank_layout)

        # 获取图片尺寸
        try:
            with Image.open(image_path) as img:
                img_width, img_height = img.size
        except Exception as e:
            print(f"无法读取图片尺寸: {e}")
            return

        pixels_to_emu = self.PIXELS_TO_EMU
        layout = self._resolve_slide_layout(img_width, img_height)
        self._ensure_slide_dimensions(layout["slide_width"], layout["slide_height"])
        if layout["ppt_scale"] != 1.0:
            print(
                f"[INFO] 图片尺寸超过PPT限制，缩放: {img_width}x{img_height} -> "
                f"{int(round(img_width * layout['ppt_scale']))}x{int(round(img_height * layout['ppt_scale']))} "
                f"(比例: {layout['ppt_scale']:.4f})"
            )

        img_left = int(layout["left"] * pixels_to_emu)
        img_top = int(layout["top"] * pixels_to_emu)
        ppt_img_width = int(layout["image_width"] * pixels_to_emu)
        ppt_img_height = int(layout["image_height"] * pixels_to_emu)

        # 添加图片
        try:
            slide.shapes.add_picture(
                image_path,
                img_left, img_top,
                width=ppt_img_width,
                height=ppt_img_height
            )
        except Exception as e:
            print(f"添加图片失败: {e}")
            return

        # 添加可编辑的文本框
        for item in text_boxes:
            rect = item.get('rect') if isinstance(item, dict) else None
            if rect is None:
                continue

            # 直接使用传入的坐标（已经在主程序中还原过了）
            x, y, w, h = self._rect_to_xywh(rect, img_width, img_height)

            content_scale = float(layout["content_scale"])
            x = int(round(float(x) * content_scale)) + int(layout["left"])
            y = int(round(float(y) * content_scale)) + int(layout["top"])
            w = max(1, int(round(float(w) * content_scale)))
            h = max(1, int(round(float(h) * content_scale)))

            text = item.get('text', "")

            # 调试：打印文本内容
            print(f"  添加文本框: {text} (位置: {x}, {y}, 大小: {w}x{h})")

            # OCR 可能会出现 0/负数宽高，跳过避免生成异常或巨大文本框
            if w <= 0 or h <= 0:
                continue

            ppt_x = int(x * pixels_to_emu)
            ppt_y = int(y * pixels_to_emu)
            # 使用原始宽度和高度，不增加余量
            ppt_w = int(w * pixels_to_emu)
            ppt_h = int(h * pixels_to_emu)

            # 避免 PowerPoint 自动调整文本框大小导致“变大”
            # 这里保留 OCR 的原始框大小（像素->EMU 1:1 映射）
            ppt_w = max(1, ppt_w)
            ppt_h = max(1, ppt_h)

            # 创建文本框
            textbox = slide.shapes.add_textbox(ppt_x, ppt_y, ppt_w, ppt_h)
            text_frame = textbox.text_frame

            # 设置文本
            text_frame.auto_size = MSO_AUTO_SIZE.NONE  # 禁用"根据文本自动调整形状大小"
            text_frame.text = str(text)  # 确保是字符串
            text_frame.word_wrap = False  # 禁用自动换行，保持单行
            text_frame.margin_left = Pt(1)
            text_frame.margin_right = Pt(1)
            text_frame.margin_top = Pt(1)
            text_frame.margin_bottom = Pt(1)
            text_frame.vertical_anchor = MSO_ANCHOR.TOP  # 顶部对齐，不使用垂直居中

            # 设置段落格式
            p = text_frame.paragraphs[0]
            # 对齐
            align = (item.get("align") if isinstance(item, dict) else "left") or "left"
            align = str(align).lower()
            if align == "center":
                p.alignment = PP_ALIGN.CENTER
            elif align == "right":
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT

            # 使用二分查找法计算最佳字体大小，确保文本完全适配在文本框内
            # 注意：这里的 w, h 已经是缩放后的像素尺寸
            font_size = None
            if isinstance(item, dict):
                fs = item.get("font_size")
                try:
                    font_size = int(fs) if fs is not None else None
                except Exception:
                    font_size = None
            if font_size is None:
                # 计算字体大小，并应用额外缩小系数
                # 因为 PPT 的文本渲染引擎与画布预览不同，实际显示会偏大
                calculated_size = self.fit_font_size(text, w, h, padding_x=2, padding_y=2)
                font_size = int(calculated_size * 0.70)
            elif layout["layout_scale"] != 1.0:
                font_size = max(1, int(round(float(font_size) * float(layout["layout_scale"]))))
            p.font.size = Pt(int(font_size))

            print(f"    文本框尺寸: {w}x{h}px, 字体: {font_size}pt")

            # 字体（默认微软雅黑）
            family = "微软雅黑"
            if isinstance(item, dict) and item.get("font_family"):
                family = str(item.get("font_family"))
            p.font.name = family
            try:
                from pptx.oxml.shared import OxmlElement
                rPr = p._element.get_or_add_rPr()
                ea = OxmlElement('a:ea')
                ea.set('typeface', family)
                rPr.append(ea)
            except:
                pass

            # 加粗
            if isinstance(item, dict):
                p.font.bold = bool(item.get("bold", False))

            # 文字颜色
            color_rgb = (0, 0, 0)
            if isinstance(item, dict):
                tc = item.get("text_color", [0, 0, 0])
                if isinstance(tc, (list, tuple)) and len(tc) == 3:
                    try:
                        color_rgb = (int(tc[0]), int(tc[1]), int(tc[2]))
                    except Exception:
                        color_rgb = (0, 0, 0)
            p.font.color.rgb = RGBColor(*color_rgb)

            print(f"    字体: {font_size}pt")

            # 设置文本框样式：支持每个文本框单独背景色（use_custom_bg + bg_color=[r,g,b]）
            box_bg_color = None
            if isinstance(item, dict):
                bg = item.get("bg_color")
                if item.get("use_custom_bg") and isinstance(bg, (list, tuple)) and len(bg) == 3:
                    try:
                        box_bg_color = (int(bg[0]), int(bg[1]), int(bg[2]))
                    except Exception:
                        box_bg_color = None

            # 单框背景色优先，其次全局背景色
            if box_bg_color is None:
                box_bg_color = self.text_bg_color

            if box_bg_color:
                textbox.fill.solid()
                r, g, b = box_bg_color
                textbox.fill.fore_color.rgb = RGBColor(r, g, b)
                # 透明度（alpha:0-255 -> transparency:0-1）
                alpha = self.text_bg_alpha
                # 只有“单框自定义背景”才使用单框 alpha；全局背景使用全局 alpha
                if isinstance(item, dict) and item.get("use_custom_bg"):
                    try:
                        alpha = int(item.get("bg_alpha", self.text_bg_alpha))
                    except Exception:
                        alpha = self.text_bg_alpha
                alpha = max(0, min(alpha, 255))
                self._set_shape_fill_alpha(textbox, alpha)
            else:
                textbox.fill.background()

            # 无边框
            # 边框
            # Border is intentionally disabled for this project.
            textbox.line.fill.background()

        print(f"[OK] 已添加页面 (图片尺寸: {img_width}x{img_height}, 共 {len(text_boxes)} 个文本框)")

    def _set_shape_fill_alpha(self, shape, alpha_0_255: int):
        """设置 shape 填充透明度（python-pptx 没有公开 API，只能改 XML）"""
        try:
            from pptx.oxml.xmlchemy import OxmlElement
            spPr = shape._element.spPr
            solidFill = spPr.find(qn("a:solidFill"))
            if solidFill is None:
                return
            # 常见是 a:srgbClr，也可能是 a:schemeClr
            clr = solidFill.find(qn("a:srgbClr"))
            if clr is None:
                clr = solidFill.find(qn("a:schemeClr"))
            if clr is None:
                return
            alpha = clr.find(qn("a:alpha"))
            if alpha is None:
                alpha = OxmlElement("a:alpha")
                clr.append(alpha)
            # DrawingML: val=0..100000（不透明度百分比*1000）
            val = int(max(0, min(255, int(alpha_0_255))) / 255.0 * 100000)
            alpha.set("val", str(val))
        except Exception:
            pass

    def fit_font_size(self, text, box_w_px, box_h_px, min_pt=6, max_pt=200, dpi=96, padding_x=6, padding_y=2):
        """
        使用二分查找法计算最佳字体大小，确保文字能完全适配文本框

        Args:
            text: 文本内容
            box_w_px: 文本框宽度（像素）
            box_h_px: 文本框高度（像素）
            min_pt: 最小字体大小（pt）
            max_pt: 最大字体大小（pt），默认200以支持大尺寸图片
            dpi: DPI（默认96）
            padding_x: 水平padding（像素）
            padding_y: 垂直padding（像素）

        Returns:
            最佳字体大小（pt）
        """
        text = (text or "").strip()
        if not text:
            return max(min_pt, min(12, max_pt))

        box_w_px = int(box_w_px or 0)
        box_h_px = int(box_h_px or 0)
        if box_w_px <= 0 or box_h_px <= 0:
            return max(min_pt, min(12, max_pt))

        avail_w = max(1, box_w_px - padding_x)
        avail_h = max(1, box_h_px - padding_y)

        # 尝试获取微软雅黑字体
        font_path = self._get_font_path("微软雅黑")
        if not font_path:
            # 如果找不到字体，使用基于高度的估算
            est = int(avail_h * 72 / dpi * 0.8)
            return max(min_pt, min(est, max_pt))

        # 创建临时画布用于测量文字
        draw = ImageDraw.Draw(Image.new("RGB", (8, 8)))

        def fits(pt):
            """测试指定字体大小是否能适配文本框"""
            px = max(1, int(round(pt * dpi / 72)))
            try:
                font = ImageFont.truetype(font_path, px)
            except:
                return True
            try:
                bbox = draw.textbbox((0, 0), text, font=font)
                w = bbox[2] - bbox[0]
                h = bbox[3] - bbox[1]
            except:
                return True
            return w <= avail_w and h <= avail_h

        # 二分查找最大的适配字体大小
        # 初始上限基于高度估算，但不要太激进
        hi = min(max_pt, max(min_pt, int(avail_h * 72 / dpi * 1.2)))
        lo = min_pt
        best = min_pt

        for _ in range(15):  # 增加迭代次数以获得更精确的结果
            mid = (lo + hi) // 2
            if fits(mid):
                best = mid
                lo = mid + 1
            else:
                hi = mid - 1

        return max(min_pt, min(best, max_pt))

    def _get_font_path(self, font_name):
        """获取字体文件路径"""
        # Windows字体目录
        font_dir = "C:/Windows/Fonts"

        # 常见字体文件名映射
        font_map = {
            "微软雅黑": ["msyh.ttc", "msyh.ttf"],
            "宋体": ["simsun.ttc", "simsun.ttf"],
            "黑体": ["simhei.ttf"],
            "Arial": ["arial.ttf"],
        }

        if font_name in font_map:
            for filename in font_map[font_name]:
                path = os.path.join(font_dir, filename)
                if os.path.exists(path):
                    return path

        return None

    def calculate_font_and_spacing(self, text, box_width, box_height):
        """
        精确计算字体大小和字符间距，完美适配文本框

        Args:
            text: 文本内容
            box_width: 文本框宽度（EMU单位）
            box_height: 文本框高度（EMU单位）

        Returns:
            (font_size, char_spacing): 字体大小（pt）和字符间距（pt）
        """
        if not text:
            return 12, 0

        # 转换为点数
        width_pt = box_width / Pt(1)
        height_pt = box_height / Pt(1)

        # 去除边距后的可用空间
        available_width = width_pt - 2  # 左右边距各1pt
        available_height = height_pt - 2  # 上下边距各1pt

        # 统计中英文字符
        chinese_count = sum(1 for c in text if '\u4e00' <= c <= '\u9fff')
        english_count = len(text) - chinese_count
        total_chars = len(text)

        if total_chars == 0:
            return 12, 0

        # 步骤1：根据高度确定字体大小上限
        # 微软雅黑字体，实际高度约为字号的0.9倍
        max_font_by_height = int(available_height / 0.9)

        # 步骤2：根据宽度计算字体大小
        # 微软雅黑：中文字符宽度 ≈ 字号 × 1.0，英文字符宽度 ≈ 字号 × 0.5
        # 初始不考虑字间距
        if chinese_count > 0 or english_count > 0:
            # 估算总宽度需求（不含字间距）
            estimated_width_per_pt = chinese_count * 1.0 + english_count * 0.5
            estimated_font_size = int(available_width / estimated_width_per_pt)
        else:
            estimated_font_size = 12

        # 步骤3：取两者中的较小值
        font_size = min(estimated_font_size, max_font_by_height)

        # 限制字体大小范围
        font_size = max(6, min(font_size, 72))

        # 步骤4：计算字符间距
        # 计算当前字体大小下文本的实际宽度
        actual_text_width = chinese_count * font_size * 1.0 + english_count * font_size * 0.5

        # 计算剩余空间
        remaining_width = available_width - actual_text_width

        # 如果有剩余空间，分配给字符间距
        if remaining_width > 0 and total_chars > 1:
            # 字符间距 = 剩余空间 / (字符数 - 1)
            char_spacing = remaining_width / (total_chars - 1)
            # 限制字间距范围（-2pt 到 10pt）
            char_spacing = max(-2, min(char_spacing, 10))
        else:
            # 如果空间不足，使用负字间距压缩
            char_spacing = remaining_width / total_chars if total_chars > 0 else 0
            char_spacing = max(-2, char_spacing)

        # 四舍五入到小数点后1位
        char_spacing = round(char_spacing, 1)

        return font_size, char_spacing

    def set_character_spacing(self, paragraph, spacing_pt):
        """
        设置字符间距

        Args:
            paragraph: 段落对象
            spacing_pt: 字符间距（磅）
        """
        try:
            from pptx.oxml.shared import OxmlElement
            from pptx.util import Pt

            # 获取段落的运行元素
            for run in paragraph.runs:
                rPr = run._element.get_or_add_rPr()

                # 设置字符间距（单位：百分之一磅）
                spacing = OxmlElement('a:spc')
                spacing.set('val', str(int(spacing_pt * 100)))
                rPr.append(spacing)
        except Exception as e:
            print(f"设置字符间距失败: {e}")

    def save(self, output_path: str):
        """
        保存PPT文件

        Args:
            output_path: 输出路径
        """
        try:
            self.prs.save(output_path)
            print(f"[OK] PPT已保存: {output_path}")
            return True
        except Exception as e:
            print(f"[X] PPT保存失败: {e}")
            return False


# 测试代码
if __name__ == "__main__":
    import cv2
    import numpy as np

    # 创建测试图片
    img = np.ones((400, 600, 3), dtype=np.uint8) * 255
    cv2.putText(img, "Test Image", (50, 200),
                cv2.FONT_HERSHEY_SIMPLEX, 2, (0, 0, 255), 3)

    test_img_path = "test_img.png"
    cv2.imwrite(test_img_path, img)

    # 创建PPT
    exporter = PPTExporter()
    exporter.add_image_with_text_boxes(
        test_img_path,
        [],
        title="测试页面"
    )
    exporter.save("test_output.pptx")

    # 清理
    if os.path.exists(test_img_path):
        os.remove(test_img_path)

    print("测试完成！")
